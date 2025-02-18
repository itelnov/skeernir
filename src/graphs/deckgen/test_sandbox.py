from deckgen.utils import run_code

test_string = (
"""
from pptx import Presentation
from pptx.util import Inches
import os

def compile_presentation():

    prs = Presentation()

    # Slide 1: Title Slide
    title_slide_layout = prs.slide_layouts[0]  # Access layout by index (Title Slide)
    title_slide = prs.slides.add_slide(title_slide_layout)

    # Add title and subtitle
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1] if len(title_slide.placeholders) > 1 else None

    title.text = "China - Introduction"
    if subtitle:
        subtitle.text = "A Glimpse into China's History, Culture, and Economy"

    # Slide 2: Content Slide
    content_slide_layout = prs.slide_layouts[1]  # Access layout by index (Content Slide)
    content_slide = prs.slides.add_slide(content_slide_layout)

    # Add title to content slide
    content_title = content_slide.shapes.title
    content_title.text = "Key Facts about China"

    # Add bullet points
    content_placeholder = content_slide.placeholders[1] if len(content_slide.placeholders) > 1 else None

    if content_placeholder:
        text_frame = content_placeholder.text_frame  # Access the text frame

        # Set the first bullet point (initial paragraph)
        text_frame.text = "- Official Name: People's Republic of China (PRC)"

        # Add new bullet points by appending paragraphs
        p = text_frame.add_paragraph()
        p.text = "- Geography: Located in East Asia, bordered by 14 countries including Russia, India, and Vietnam."

        p = text_frame.add_paragraph()
        p.text = "- Population: Over 1.4 billion people, making it the most populous country in the world."

        p = text_frame.add_paragraph()
        p.text = "- Capital City: Beijing"

        p = text_frame.add_paragraph()
        p.text = "- Language: Mandarin is the official language, with various dialects spoken across regions."

        p = text_frame.add_paragraph()
        p.text = "- Quick Fact: China has a rich history spanning over 5,000 years and is home to one of the Seven Wonders of the Ancient World, the Great Wall of China."

    # Save presentation to a file
    prs.save(os.path.join(".outputs", "China_Presentation.pptx"))

# Execute the function
compile_presentation()
""")

if __name__ == '__main__':
    result, stdout, stderr = run_code(
        test_string, allowed_path=".outputs")

    print(f"Result: {result}")
    print(f"Output: {stdout}")
    print(f"Errors: {stderr}")